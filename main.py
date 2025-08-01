import os
from flask import Flask, request, jsonify, send_file
from flask_sqlalchemy import SQLAlchemy
from werkzeug.security import generate_password_hash, check_password_hash
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
import io
import base64 # Required for handling image data

app = Flask(__name__)

# --- DATABASE CONFIGURATION ---
app.config['SQLALCHEMY_DATABASE_URI'] = os.environ.get('DATABASE_URL')
app.config['SQLALCHEMY_TRACK_MODIFICATIONS'] = False
db = SQLAlchemy(app)

# --- DATABASE MODELS ---
class User(db.Model):
    id = db.Column(db.Integer, primary_key=True)
    username = db.Column(db.String(80), unique=True, nullable=False)
    password_hash = db.Column(db.String(128), nullable=False)

class Song(db.Model):
    id = db.Column(db.Integer, primary_key=True)
    title = db.Column(db.String(120), nullable=False)
    lyrics = db.Column(db.Text, nullable=False)
    user_id = db.Column(db.Integer, db.ForeignKey('user.id'), nullable=False)

# --- COMMAND TO CREATE DATABASE TABLES ---
@app.cli.command("create-db")
def create_db():
    with app.app_context():
        db.create_all()
    print("Database tables created!")

# --- API ENDPOINTS (Unchanged) ---
@app.route('/register', methods=['POST'])
def register():
    data = request.get_json()
    if not data or 'username' not in data or 'password' not in data:
        return jsonify({'message': 'Username and password are required!'}), 400
    if User.query.filter_by(username=data['username']).first():
        return jsonify({'message': 'Username already exists!'}), 409
    hashed_password = generate_password_hash(data['password'], method='pbkdf2:sha256')
    new_user = User(username=data['username'], password_hash=hashed_password)
    db.session.add(new_user)
    db.session.commit()
    return jsonify({'message': 'New user created!'}), 201

@app.route('/login', methods=['POST'])
def login():
    data = request.get_json()
    if not data or 'username' not in data or 'password' not in data:
        return jsonify({'message': 'Could not verify'}), 401
    user = User.query.filter_by(username=data['username']).first()
    if not user or not check_password_hash(user.password_hash, data['password']):
        return jsonify({'message': 'Could not verify! Incorrect username or password.'}), 401
    return jsonify({'message': 'Login successful!', 'user_id': user.id}), 200

@app.route('/songs', methods=['POST'])
def add_song():
    data = request.get_json()
    if not data or 'title' not in data or 'lyrics' not in data or 'user_id' not in data:
        return jsonify({'message': 'Missing data!'}), 400
    new_song = Song(title=data['title'], lyrics=data['lyrics'], user_id=data['user_id'])
    db.session.add(new_song)
    db.session.commit()
    return jsonify({'message': 'Song created!', 'song_id': new_song.id}), 201

@app.route('/songs/<int:user_id>', methods=['GET'])
def get_songs(user_id):
    songs = Song.query.filter_by(user_id=user_id).order_by(Song.title).all()
    output = []
    for song in songs:
        song_data = {'id': song.id, 'title': song.title, 'lyrics': song.lyrics}
        output.append(song_data)
    return jsonify({'songs': output})

@app.route('/songs/<int:song_id>', methods=['PUT', 'DELETE'])
def manage_song(song_id):
    song = Song.query.get(song_id)
    if not song:
        return jsonify({'message': 'Song not found!'}), 404
    if request.method == 'PUT':
        data = request.get_json()
        song.title = data.get('title', song.title)
        song.lyrics = data.get('lyrics', song.lyrics)
        db.session.commit()
        return jsonify({'message': 'Song updated!'})
    if request.method == 'DELETE':
        db.session.delete(song)
        db.session.commit()
        return jsonify({'message': 'Song deleted!'})

# --- HELPER FUNCTION ---
def hex_to_rgb(hex_color, default_color=(0, 0, 0)):
    hex_color = hex_color.lstrip('#').strip()
    if not hex_color: return default_color
    if len(hex_color) == 3: hex_color = "".join([c*2 for c in hex_color])
    if len(hex_color) == 6:
        try:
            return tuple(int(hex_color[i:i+2], 16) for i in (0, 2, 4))
        except ValueError:
            return default_color
    return default_color

# --- POWERPOINT GENERATION ENDPOINT (MODIFIED) ---
@app.route('/generate-ppt', methods=['POST'])
def generate_ppt_custom():
    try:
        data = request.get_json()
        lyrics = data.get('lyrics', '')
        song_title = data.get('title', 'Lyrics')
        bg_color_hex = data.get('backgroundColor', '000000')
        font_color_hex = data.get('fontColor', 'FFFFFF')
        font_size = int(data.get('fontSize', 44))
        
        # ✅ Get the new style properties
        font_name = data.get('fontName', 'Arial') 
        background_image_b64 = data.get('backgroundImage')

        prs = Presentation()
        prs.slide_width = Inches(16)
        prs.slide_height = Inches(9)
        blank_slide_layout = prs.slide_layouts[6]
        
        # Since we send one slide at a time, we just use the lyrics directly
        slide_content = lyrics.strip()

        if not slide_content:
             return "Cannot generate an empty presentation.", 400
        
        slide = prs.slides.add_slide(blank_slide_layout)
        
        # ✅ New Background Logic: Prioritize Image over Color
        if background_image_b64:
            try:
                image_data = base64.b64decode(background_image_b64)
                image_stream = io.BytesIO(image_data)
                slide.background.fill.picture(image_stream)
            except Exception as e:
                print(f"Error decoding or using background image: {e}")
                # Fallback to color if image fails
                fill = slide.background.fill
                fill.solid()
                fill.fore_color.rgb = RGBColor(*hex_to_rgb(bg_color_hex, (0,0,0)))
        else:
            # Fallback to color if no image is provided
            fill = slide.background.fill
            fill.solid()
            fill.fore_color.rgb = RGBColor(*hex_to_rgb(bg_color_hex, (0,0,0)))

        # Add text box and style it
        txBox = slide.shapes.add_textbox(Inches(1.0), Inches(1.0), Inches(14.0), Inches(7.0))
        tf = txBox.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = slide_content
        p.alignment = PP_ALIGN.CENTER
        
        font = p.font
        font.name = font_name # ✅ Apply font family
        font.size = Pt(font_size)
        font.color.rgb = RGBColor(*hex_to_rgb(font_color_hex, (255,255,255)))

        # Save to a memory stream and send the file
        file_stream = io.BytesIO()
        prs.save(file_stream)
        file_stream.seek(0)
        
        safe_filename = "".join([c for c in song_title if c.isalpha() or c.isdigit() or c==' ']).rstrip()
        return send_file(
            file_stream,
            as_attachment=True,
            download_name=f"{safe_filename}.pptx",
            mimetype='application/vnd.openxmlformats-officedocument.presentationml.presentation'
        )
    except Exception as e:
        print(f"Error in /generate-ppt: {e}")
        return jsonify({"error": str(e)}), 500